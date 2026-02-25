from __future__ import annotations
from dataclasses import dataclass
from pathlib import Path
from typing import List
import re

@dataclass
class SlideText:
    slide_id: str
    text: str

# ---------------------------
# Cleaning
# ---------------------------

def _clean(t: str) -> str:
    if not t:
        return ""

    t = t.replace("\u00a0", " ")
    t = re.sub(r"\s+", " ", t).strip()

    # remove garbage short tokens
    words = [w for w in t.split() if len(w) > 2]
    t = " ".join(words)

    return t


# ---------------------------
# PPTX extraction (IMPROVED)
# ---------------------------

def extract_from_pptx(pptx_path: str | Path) -> List[SlideText]:
    from pptx import Presentation

    prs = Presentation(str(pptx_path))
    out: List[SlideText] = []

    for i, slide in enumerate(prs.slides, start=1):
        chunks = []

        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                txt = shape.text.strip()
                if len(txt) > 3:
                    chunks.append(txt)

        text = _clean(" ".join(chunks))

        out.append(SlideText(
            slide_id=f"slide_{i}",
            text=text if text else "(empty slide)"
        ))

    return out


# ---------------------------
# PDF extraction (TEXT)
# ---------------------------

def extract_from_pdf_text(pdf_path: str | Path) -> List[SlideText]:
    import pdfplumber

    out: List[SlideText] = []

    with pdfplumber.open(str(pdf_path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            text = page.extract_text() or ""
            text = _clean(text)

            out.append(SlideText(
                slide_id=f"page_{i}",
                text=text if text else "(empty page)"
            ))

    return out


# ---------------------------
# OCR fallback (STRONGER)
# ---------------------------

def ocr_pdf_pages(pdf_path: str | Path) -> List[SlideText]:
    import pdfplumber
    import pytesseract

    out: List[SlideText] = []

    with pdfplumber.open(str(pdf_path)) as pdf:
        for i, page in enumerate(pdf.pages, start=1):
            try:
                im = page.to_image(resolution=300).original
                text = pytesseract.image_to_string(im)
            except Exception:
                text = ""

            text = _clean(text)

            out.append(SlideText(
                slide_id=f"page_{i}",
                text=text if text else "(ocr failed)"
            ))

    return out


# ---------------------------
# MAIN LOADER
# ---------------------------

def load_slides(file_path: str | Path, use_ocr_if_empty: bool = False) -> List[SlideText]:
    p = Path(file_path)

    if not p.exists():
        raise FileNotFoundError(p)

    if p.suffix.lower() == ".pptx":
        slides = extract_from_pptx(p)

    elif p.suffix.lower() == ".pdf":
        slides = extract_from_pdf_text(p)

        # If too many empty pages → OCR
        empty_ratio = sum(1 for s in slides if s.text == "(empty page)") / max(1, len(slides))

        if use_ocr_if_empty and empty_ratio > 0.4:
            print("⚠️ Falling back to OCR for PDF...")
            slides = ocr_pdf_pages(p)

    else:
        raise ValueError("Only .pptx or .pdf supported")

    return slides